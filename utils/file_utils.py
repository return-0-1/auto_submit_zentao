import json
import logging
import re
from typing import Dict, Any, Optional
import win32api

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


def read_txt_file(file_path: str) -> str:
    """读取txt文件内容"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except FileNotFoundError:
        logging.error(f"文件不存在: {file_path}")
        return ""
    except Exception as e:
        logging.error(f"读取文件失败: {file_path}, 错误: {str(e)}")
        return ""


def read_json_file(file_path: str) -> Dict[str, Any]:
    """安全读取JSON文件，包含异常处理"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return json.load(file)
    except FileNotFoundError:
        logging.error(f"JSON文件不存在: {file_path}")
        raise
    except json.JSONDecodeError:
        logging.error(f"JSON文件格式错误: {file_path}")
        raise


def get_url_by_type(submit_type: str) -> Optional[str]:
    """根据提交类型（case/bug）返回对应URL"""
    from config import CASE_URL, BUG_URL
    url_map = {
        "case": CASE_URL,
        "bug": BUG_URL
    }
    url = url_map.get(submit_type)
    if not url:
        logging.warning(f"无效的提交类型: {submit_type}，仅支持case/bug")
    return url


def clean_filename(raw_filename):
    """
    清理文件名，移除括号中的文件大小信息和其他多余内容

    Args:
        raw_filename: 原始文件名，如 "FeatureOutrectangleTouchFeature_新加算法需求.docx (39.87K)"

    Returns:
        str: 清理后的文件名，如 "FeatureOutrectangleTouchFeature_新加算法需求.docx"
    """
    cleaned = re.sub(r'\s*\([^)]*\)\s*', '', raw_filename)

    if cleaned == raw_filename:
        cleaned = re.sub(r'\s*\d*\.?\d+[KMGTPE]?B?\s*$', '', raw_filename)

    if cleaned == raw_filename:
        cleaned = re.sub(r'\s+[\d\.]+\s*[KMGTP]?B?\s*$', '', raw_filename)

    if not cleaned.strip():
        cleaned = raw_filename

    return cleaned.strip()


def clean_filename_advanced(raw_filename):
    """
    高级文件名清理方法，处理更多情况
    """
    patterns = [
        r'\s*\(\d+\.\d+[KMGTP]?\)\s*',
        r'\s*\d+\.\d+[KMGTP]?\s*',
        r'\s*\(\d+[KMGTP]?\)\s*',
        r'\s*\d+[KMGTP]?\s*',
        r'\s*\[[^\]]*\]\s*',
        r'\s*【[^】]*】\s*',
    ]

    cleaned = raw_filename
    for pattern in patterns:
        cleaned = re.sub(pattern, '', cleaned)

    cleaned = cleaned.strip().strip('.')

    return cleaned if cleaned else raw_filename


# ===== PPTX/DOCX 文本提取相关函数 =====

def extract_text_from_pptx(pptx_path: str) -> str:
    """从单个PPTX文件中提取文本"""
    from pptx import Presentation
    
    prs = Presentation(pptx_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text.append(cell.text)

    return '\n'.join(text)


def extract_text_from_docx(docx_path: str) -> str:
    """从DOCX文件中提取文本内容"""
    try:
        try:
            import docx
            doc = docx.Document(docx_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"

            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + "\t"
                    text += "\n"
                text += "\n"

            return text.strip()

        except ImportError:
            import zipfile
            import xml.etree.ElementTree as ET

            text = ""
            with zipfile.ZipFile(docx_path) as docx_zip:
                if 'word/document.xml' in docx_zip.namelist():
                    xml_content = docx_zip.read('word/document.xml')
                    root = ET.fromstring(xml_content)

                    namespaces = {
                        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                    }

                    for elem in root.iter():
                        if elem.tag.endswith('t'):
                            if elem.text:
                                text += elem.text
                        elif elem.tag.endswith('tab'):
                            text += '\t'
                        elif elem.tag.endswith('br'):
                            text += '\n'
                        elif elem.tag.endswith('p'):
                            text += '\n'

            return text.strip()

    except Exception as e:
        raise Exception(f"提取DOCX文本失败: {str(e)}")


def process_single_story_folder(story_folder_path: str, story_id: str, output_base_folder: str):
    """处理单个需求文件夹中的所有PPTX和DOCX文件"""
    import os
    
    all_text_content = []

    for filename in os.listdir(story_folder_path):
        filepath = os.path.join(story_folder_path, filename)

        if filename.endswith('.pptx') or filename.endswith('.ppt'):
            try:
                logging.info(f"  正在处理PPT文件: {filename}")
                text = extract_text_from_pptx(filepath)
                file_content = f"=== 文件: {filename} ===\n{text}\n\n"
                all_text_content.append(file_content)
                logging.info(f"  成功提取PPT: {filename}")

            except Exception as e:
                error_msg = f"处理PPT文件 {filename} 时出错: {str(e)}"
                logging.error(f"  {error_msg}")
                all_text_content.append(f"=== 文件: {filename} [处理失败] ===\n{error_msg}\n\n")

        elif filename.endswith('.docx'):
            try:
                logging.info(f"  正在处理Word文件: {filename}")
                text = extract_text_from_docx(filepath)
                file_content = f"=== 文件: {filename} ===\n{text}\n\n"
                all_text_content.append(file_content)
                logging.info(f"  成功提取Word: {filename}")

            except Exception as e:
                error_msg = f"处理Word文件 {filename} 时出错: {str(e)}"
                logging.error(f"  {error_msg}")
                all_text_content.append(f"=== 文件: {filename} [处理失败] ===\n{error_msg}\n\n")

    if all_text_content:
        from pathlib import Path
        
        output_filename = f"{story_id}.txt"
        output_path = os.path.join(output_base_folder, output_filename)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"=== 需求ID: {story_id} ===\n")
            f.write(f"=== 源文件夹: {os.path.basename(story_folder_path)} ===\n\n")

            for content in all_text_content:
                f.write(content)

        logging.info(f"✓ 需求 {story_id} 的文本已保存到: {output_path}")
    else:
        logging.warning(f"⚠ 需求 {story_id} 文件夹中没有找到PPTX或DOCX文件")


def extract_text_from_story_folders(root_folder: str, output_base_folder: str = '零一'):
    """从story_xxx格式的文件夹中提取PPTX文本"""
    import os
    from pathlib import Path
    
    Path(output_base_folder).mkdir(exist_ok=True)

    for item in os.listdir(root_folder):
        item_path = os.path.join(root_folder, item)

        if os.path.isdir(item_path) and item.startswith('story_'):
            try:
                story_id = item.replace('story_', '')
                logging.info(f"正在处理需求文件夹: {item} (ID: {story_id})")
                process_single_story_folder(item_path, story_id, output_base_folder)

            except Exception as e:
                logging.error(f"处理文件夹 {item} 时出错: {str(e)}")


def extract_text_from_multiple_pptx(folder_path: str, output_folder: str = 'output'):
    """从文件夹中的多个PPTX文件提取文本"""
    import os
    from pathlib import Path
    
    Path(output_folder).mkdir(exist_ok=True)

    for filename in os.listdir(folder_path):
        if filename.endswith('.pptx') or filename.endswith('.ppt'):
            try:
                filepath = os.path.join(folder_path, filename)
                logging.info(f"正在处理文件: {filename}")

                text = extract_text_from_pptx(filepath)

                output_filename = os.path.splitext(filename)[0] + '.txt'
                output_path = os.path.join(output_folder, output_filename)

                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)

                logging.info(f"已保存: {output_path}")

            except Exception as e:
                logging.error(f"处理文件 {filename} 时出错: {str(e)}")


def merge_txt_files(input_folder: str, output_file: str = 'merged_documents.txt'):
    """合并文件夹中的所有TXT文档"""
    import os
    from pathlib import Path
    
    if not os.path.isdir(input_folder):
        raise ValueError(f"输入文件夹不存在: {input_folder}")

    txt_files = sorted([f for f in os.listdir(input_folder) if f.lower().endswith('.txt')])

    if not txt_files:
        logging.warning(f"警告: 文件夹 {input_folder} 中没有找到TXT文件")
        return

    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with open(output_file, 'w', encoding='utf-8') as outfile:
        for filename in txt_files:
            filepath = os.path.join(input_folder, filename)

            outfile.write(f"\n\n=== 文档: {filename} ===\n\n")

            try:
                with open(filepath, 'r', encoding='utf-8') as infile:
                    content = infile.read()
                    outfile.write(content)
            except UnicodeDecodeError:
                try:
                    with open(filepath, 'r', encoding='gbk') as infile:
                        content = infile.read()
                        outfile.write(content)
                except Exception as e:
                    logging.error(f"无法读取文件 {filename}: {str(e)}")
                    outfile.write(f"[无法读取此文件内容: {str(e)}]")
            except Exception as e:
                logging.error(f"处理文件 {filename} 时出错: {str(e)}")
                outfile.write(f"[处理此文件时出错: {str(e)}]")

    logging.info(f"成功合并 {len(txt_files)} 个文档到 {output_file}")
